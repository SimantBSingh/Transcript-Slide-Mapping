#!/usr/bin/env python3
import os
import re
import io
import pickle
from dataclasses import dataclass
from typing import List, Optional
from collections import defaultdict
from datetime import datetime

import numpy as np
from sentence_transformers import SentenceTransformer
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import nltk
from nltk.stem import PorterStemmer
from nltk.corpus import stopwords

# --------------------
# Data Models
# --------------------
@dataclass
class Utterance:
    speaker: str
    text: str
    start_sec: float
    end_sec: float
    mid_sec: float
    duration: float
    utterance_id: Optional[int] = None
    slide_number: int = 0
    confidence: float = 0.0

@dataclass
class Slide:
    slide_number: int
    title: str
    raw_text: str
    content: str

# --------------------
# Setup
# --------------------
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
stemmer = PorterStemmer()
STOP_WORDS = set(stopwords.words('english'))
# Load a BERT-based sentence transformer
embedder = SentenceTransformer('bert-base-nli-mean-tokens')

# --------------------
# Utility Functions
# --------------------
def preprocess_text(text: str) -> str:
    # Simple cleanup: remove extra whitespace
    return re.sub(r"\s+", " ", text.strip())

def parse_timestamp(timestamp_str: str) -> dict:
    m = re.findall(r"(\d{2}):(\d{2}):(\d{2})\.(\d{3})", timestamp_str)
    if len(m) != 2:
        return {'start_sec': 0.0, 'end_sec': 0.0, 'mid_sec': 0.0, 'duration': 0.0}
    def to_sec(h, m, s, ms): return int(h)*3600 + int(m)*60 + int(s) + int(ms)/1000
    start = to_sec(*m[0]); end = to_sec(*m[1])
    return {'start_sec': start, 'end_sec': end, 'mid_sec': (start+end)/2, 'duration': end-start}

# --------------------
# Extract Slides
# --------------------
def extract_slides(pptx_path: str) -> List[Slide]:
    prs = Presentation(pptx_path)
    slides: List[Slide] = []
    for idx, sl in enumerate(prs.slides, start=1):
        title = ''
        placeholders = [s for s in sl.shapes if s.is_placeholder and s.placeholder_format.type == PP_PLACEHOLDER.TITLE]
        if placeholders:
            shp = placeholders[0]
            if shp.has_text_frame:
                title = shp.text.strip()
        if not title:
            for s in sl.shapes:
                if s.has_text_frame and s.text.strip():
                    title = s.text.strip()
                    break
        contents = []
        for s in sl.shapes:
            if s.has_text_frame and (not s.is_placeholder or s.placeholder_format.type != PP_PLACEHOLDER.TITLE):
                for p in s.text_frame.paragraphs:
                    t = p.text.strip()
                    if t and t != title:
                        contents.append(t)
        rawtxt = f"{title} " + " ".join(contents)
        slides.append(Slide(idx, title, rawtxt, "\n".join(contents)))
    return slides

# --------------------
# Extract Transcript
# --------------------
def extract_raw_transcript(vtt_path: str) -> List[Utterance]:
    lines = [l.strip() for l in open(vtt_path, encoding='utf-8')]
    raw: List[Utterance] = []
    curr_ts = None
    curr_spk = "Unknown"
    buf: List[str] = []
    for ln in lines:
        if '-->' in ln:
            curr_ts = parse_timestamp(ln)
            continue
        m = re.match(r'^([A-Za-z ]+):\s(.+)$', ln)
        if m:
            if buf and curr_ts:
                text = " ".join(buf)
                text = re.sub(r'\s*\d+\s*$', '', text)
                raw.append(Utterance(curr_spk, text,
                                     curr_ts['start_sec'], curr_ts['end_sec'],
                                     curr_ts['mid_sec'], curr_ts['duration']))
                buf = []
            curr_spk = m.group(1).strip()
            buf.append(m.group(2).strip())
        elif curr_ts:
            buf.append(ln)
    if buf and curr_ts:
        text = " ".join(buf)
        text = re.sub(r'\s*\d+\s*$', '', text)
        raw.append(Utterance(curr_spk, text,
                             curr_ts['start_sec'], curr_ts['end_sec'],
                             curr_ts['mid_sec'], curr_ts['duration']))
    raw.sort(key=lambda u: u.start_sec)
    for i, u in enumerate(raw, start=1):
        u.utterance_id = i
    return raw

# --------------------
# BERT Embedding Matching
# --------------------
def encode_slides(slides: List[Slide]) -> np.ndarray:
    texts = [preprocess_text(s.raw_text) for s in slides]
    return embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=True)

def encode_utterances(utts: List[Utterance]) -> np.ndarray:
    texts = [preprocess_text(u.text) for u in utts]
    return embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=True)

def match_utterances_to_slides(utts: List[Utterance], slides: List[Slide], slide_embs: np.ndarray, utt_embs: np.ndarray) -> List[Utterance]:
    sims = cosine_similarity(utt_embs, slide_embs)
    for i, u in enumerate(utts):
        best_idx = int(np.argmax(sims[i]))
        u.slide_number = slides[best_idx].slide_number
        u.confidence = float(sims[i][best_idx])
    return utts

# --------------------
# Output Generation
# --------------------
def format_time(seconds: float) -> str:
    return datetime.utcfromtimestamp(seconds).strftime('%H:%M:%S')

def generate_transcript_by_slide(matches: List[Utterance], slides: List[Slide], out_path: str):
    groups = defaultdict(list)
    for m in matches:
        groups[m.slide_number].append(m)
    meta = {s.slide_number: s for s in slides}
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write("TRANSCRIPT BY SLIDE\n==================\n\n")
        for num in sorted(groups):
            if num == 0:
                continue
            title = meta[num].title or f"Slide {num}"
            f.write(f"SLIDE {num}: {title}\n")
            f.write('-'*80 + "\n")
            for u in groups[num]:
                f.write(f"[{format_time(u.start_sec)}] {u.speaker}: {u.text}\n")
            f.write("\n")

# --------------------
# Main Entry
# --------------------
def main(transcript_path: str, pptx_path: str, output_path: str):
    slides = extract_slides(pptx_path)
    utts = extract_raw_transcript(transcript_path)
    # Encode texts
    slide_embs = encode_slides(slides)
    utt_embs = encode_utterances(utts)
    # Match
    matches = match_utterances_to_slides(utts, slides, slide_embs, utt_embs)
    # Generate
    generate_transcript_by_slide(matches, slides, output_path)
    print(f"Generated aligned transcript: {output_path}")

if __name__ == '__main__':
    # Replace with actual paths
    main('transcript.txt', 'Lecture01_Example.pptx', 'v4_output.txt')
