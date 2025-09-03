#!/usr/bin/env python3
import os
import re
import io
import pickle
from dataclasses import dataclass
from typing import List, Optional
from collections import defaultdict, deque
from datetime import datetime

import numpy as np
from sentence_transformers import SentenceTransformer
import faiss
from kneed import KneeLocator
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import pytesseract
from PIL import Image
import nltk
from nltk.stem import PorterStemmer
from nltk.corpus import stopwords

# --------------------
# Data Models
# --------------------
@dataclass
class Utterance:
    speaker: str
    text: str
    start_sec: float
    end_sec: float
    mid_sec: float
    duration: float
    utterance_id: Optional[int] = None
    slide_number: int = 0
    confidence: float = 0.0
    confidence_type: str = ""

@dataclass
class Slide:
    slide_number: int
    title: str
    raw_text: str
    content: str

# --------------------
# One-time NLP & Embedding Setup
# --------------------
nltk.download('stopwords', quiet=True)
stemmer = PorterStemmer()
STOP_WORDS = set(stopwords.words('english'))
embedder = SentenceTransformer('all-MiniLM-L6-v2')

# --------------------
# Utility Functions
# --------------------
def parse_timestamp(timestamp_str: str) -> dict:
    pattern = r"(\d{2}):(\d{2}):(\d{2})\.(\d{3})"
    matches = re.findall(pattern, timestamp_str)
    if len(matches) != 2:
        return {'start_sec': 0.0, 'end_sec': 0.0, 'mid_sec': 0.0, 'duration': 0.0}

    def to_sec(h, m, s, ms):
        return int(h)*3600 + int(m)*60 + int(s) + int(ms)/1000

    start = to_sec(*matches[0])
    end = to_sec(*matches[1])
    return {'start_sec': start, 'end_sec': end, 'mid_sec': (start+end)/2, 'duration': end-start}

# --------------------
# Transcript Extraction
# --------------------
def extract_raw_transcript(vtt_path: str) -> List[Utterance]:
    lines = [ln.strip() for ln in open(vtt_path, encoding='utf-8')]
    raw: List[Utterance] = []
    curr_ts = None
    curr_spk = "Unknown"
    buf: List[str] = []

    for ln in lines:
        if not ln:
            continue
        if '-->' in ln:
            curr_ts = parse_timestamp(ln)
            continue
        # speaker tag or style
        m = re.match(r'^<v([^>]+)>(.+)$', ln)
        if not m:
            m = re.match(r'^([A-Z][\w\s]+):\s(.+)$', ln)
        if m:
            if buf and curr_ts:
                txt = ' '.join(buf).strip()
                # Remove trailing numeric codes
                txt = re.sub(r'\s*\d+\s*$', '', txt)
                raw.append(Utterance(curr_spk, txt,
                                     curr_ts['start_sec'], curr_ts['end_sec'],
                                     curr_ts['mid_sec'], curr_ts['duration']))
                buf.clear()
            curr_spk = m.group(1).strip()
            buf.append(m.group(2).strip())
        elif curr_ts:
            buf.append(ln)
    # last buffer
    if buf and curr_ts:
        txt = ' '.join(buf).strip()
        txt = re.sub(r'\s*\d+\s*$', '', txt)
        raw.append(Utterance(curr_spk, txt,
                             curr_ts['start_sec'], curr_ts['end_sec'],
                             curr_ts['mid_sec'], curr_ts['duration']))
    # sort & assign IDs
    raw.sort(key=lambda u: u.start_sec)
    for i, u in enumerate(raw, start=1): u.utterance_id = i
    return raw

# --------------------
# Slide Extraction
# --------------------
def extract_slides(pptx_path: str) -> List[Slide]:
    prs = Presentation(pptx_path)
    slides: List[Slide] = []
    for idx, sl in enumerate(prs.slides, start=1):
        title = ''
        placeholders = [s for s in sl.shapes if s.is_placeholder and s.placeholder_format.type == PP_PLACEHOLDER.TITLE]
        if placeholders:
            shp = placeholders[0]
            title = shp.text.strip() if shp.has_text_frame else ''
        if not title:
            for s in sl.shapes:
                if s.has_text_frame and s.text.strip():
                    title = s.text.strip(); break
        contents: List[str] = []
        for s in sl.shapes:
            if s.has_text_frame and (not s.is_placeholder or s.placeholder_format.type != PP_PLACEHOLDER.TITLE):
                for p in s.text_frame.paragraphs:
                    t = p.text.strip()
                    if t and t != title:
                        contents.append(t)
            elif not s.has_text_frame and hasattr(s, 'image'):
                # OCR fallback
                blob = s.image.blob
                try:
                    img = Image.open(io.BytesIO(blob))
                    txt = pytesseract.image_to_string(img).strip()
                    if txt: contents.append(txt)
                except Exception:
                    pass
        rawtxt = title + ' ' + ' '.join(contents)
        slides.append(Slide(idx, title, rawtxt, '\n'.join(contents)))
    return slides

# --------------------
# Precompute & Cache Slide Index
# --------------------
def build_slide_index(slides: List[Slide], index_path: str, meta_path: str):
    texts = [(s.title + ' ')*3 + s.raw_text for s in slides]
    embs = embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=True)
    dim = embs.shape[1]
    idx = faiss.IndexFlatIP(dim)
    idx.add(embs)
    with open(meta_path, 'wb') as f: pickle.dump(slides, f)
    faiss.write_index(idx, index_path)
    print(f"Built index-> {index_path}, meta-> {meta_path}")

# --------------------
# Matching & Post-Processing
# --------------------
def calculate_dynamic_threshold(scores: np.ndarray) -> float:
    nz = scores[scores>0]
    if not len(nz): return 0.0
    sorted_scores = np.sort(nz)
    kl = KneeLocator(np.arange(len(sorted_scores)), sorted_scores, curve='concave', direction='increasing')
    return float(sorted_scores[kl.knee]) if kl.knee is not None else np.percentile(nz, 80)

def match_transcript_to_slides(utts: List[Utterance], index_path: str, meta_path: str,
                             base_threshold: float=0.15) -> List[Utterance]:
    idx = faiss.read_index(index_path)
    with open(meta_path,'rb') as f: slides = pickle.load(f)
    texts = [u.text for u in utts]
    embs = embedder.encode(texts, convert_to_numpy=True, normalize_embeddings=True)
    D, I = idx.search(embs, k=1)
    scores = D[:,0]
    thr = max(base_threshold, calculate_dynamic_threshold(scores))
    for u, sc, sl_idx in zip(utts, scores, I[:,0]):
        if sc >= thr:
            u.slide_number = slides[sl_idx].slide_number
            u.confidence_type = 'high'
        else:
            u.slide_number = 0
            u.confidence_type = 'low'
        u.confidence = float(sc)
        
    closing_phrases = {"see you", "we'll continue", "next time", "wrap up"}
    for u in utts:
        if any(phrase in u.text.lower() for phrase in closing_phrases):
            if u.slide_number < len(slides) // 2:  # If matched to early slide
                u.slide_number = len(slides)  # Assign to last slide
                u.confidence_type = 'contextual'
    return utts

def temporal_analysis(matches: List[Utterance], slides: List[Slide], window_size: int=20, time_window: float=30.0):
    seq = [s.slide_number for s in slides]
    current_slide = matches[0].slide_number or seq[0]
    current_time = matches[0].start_sec
    dq = deque(maxlen=window_size)
    
    # Give more weight to utterances that occur near the end of the lecture
    total_duration = max(m.end_sec for m in matches)
    for m in matches:
        if m.start_sec > total_duration * 0.8:  # Last 20% of lecture
            m.confidence *= 1.5  # Boost confidence for closing statements
        if m.slide_number==0:
            dt = abs(m.start_sec-current_time)
            nearby = [x for x in matches if x.slide_number>0 and abs(x.start_sec-m.start_sec)<=time_window]
            if nearby:
                best = max(nearby, key=lambda x: x.confidence)
                m.slide_number, m.confidence_type = best.slide_number, 'context'
            elif dt<time_window:
                m.slide_number, m.confidence_type = current_slide, 'temporal'
            else:
                m.slide_number, m.confidence_type = current_slide, 'sequence'
        current_slide, current_time = m.slide_number, m.start_sec
        dq.append(m.slide_number)
        if len(dq)>=3:
            current_slide = max(set(dq), key=dq.count)
    return matches

def validate_matches(matches: List[Utterance], slides: List[Slide]) -> List[Utterance]:
    nums = {s.slide_number for s in slides}
    for m in matches:
        if m.slide_number not in nums:
            closest = min(nums, key=lambda x: abs(x-m.slide_number))
            m.slide_number, m.confidence_type = closest, 'corrected'
    return matches

# --------------------
# Output Generation
# --------------------
def format_time(seconds: float) -> str:
    return datetime.utcfromtimestamp(seconds).strftime('%H:%M:%S')

def generate_transcript_by_slide(matches: List[Utterance], slides: List[Slide], out_path: str):
    groups = defaultdict(list)
    for m in matches: groups[m.slide_number].append(m)
    meta = {s.slide_number:s for s in slides}
    
    with open(out_path, 'w', encoding='utf-8') as f:
        first_slide = True
        for num in sorted(groups):
            if num == 0: continue
                        
            # Add delimiter between slides, but not before the first slide
            if not first_slide:
                f.write('--' + '\n')
            else:
                first_slide = False
                
            for u in groups[num]:
                f.write(f"{u.text}\n")
            f.write('\n')

# --------------------
# Main Pipeline Entrypoint
# --------------------
def main(transcript_path: str, pptx_path: str, output_path: str):
    slides = extract_slides(pptx_path)
    index_file = 'slide_index.faiss'
    meta_file = 'slide_meta.pkl'
    if not os.path.exists(index_file) or not os.path.exists(meta_file):
        build_slide_index(slides, index_file, meta_file)

    utts = extract_raw_transcript(transcript_path)
    matches = match_transcript_to_slides(utts, index_file, meta_file, base_threshold=0.25)
    matches = temporal_analysis(matches, slides)
    matches = validate_matches(matches, slides)
    generate_transcript_by_slide(matches, slides, output_path)
    print(f"Generated aligned transcript: {output_path}")

if __name__ == '__main__':
    main('transcript.txt', 'Lecture01_Example.pptx', 'v2_output.txt')
