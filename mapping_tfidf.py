
import os
import re
import json
import numpy as np
from pptx import Presentation
from sklearn.feature_extraction.text import TfidfVectorizer
from sklearn.metrics.pairwise import cosine_similarity
import PyPDF2

def extract_slides_from_pptx(pptx_path):
    """Extract slide title and content from a PowerPoint file."""
    prs = Presentation(pptx_path)
    slides = []
    
    for slide_num, slide in enumerate(prs.slides, 1):
        title = ""
        content = []
        
        # Extract title (first shape with title text)
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame and hasattr(shape, "is_title") and shape.is_title:
                title = shape.text.strip()
                break
        
        # Extract content (non-title text)
        for shape in slide.shapes:
            if hasattr(shape, "has_text_frame") and shape.has_text_frame:
                if hasattr(shape, "is_title") and shape.is_title:
                    continue  # Skip title text
                text = shape.text.strip()
                if text:
                    content.append(text)
        
        slides.append({
            "slide_number": slide_num,
            "title": title,
            "content": " ".join(content)
        })
    
    return slides

def extract_slides_from_pdf(pdf_path):
    """Extract text content from each page of a PDF file."""
    slides = []
    
    try:
        with open(pdf_path, 'rb') as file:
            reader = PyPDF2.PdfReader(file)
            
            for page_num, page in enumerate(reader.pages, 1):
                text = page.extract_text().strip()
                
                # Simple heuristic to separate title from content
                lines = text.split('\n')
                title = lines[0] if lines else ""
                content = ' '.join(lines[1:]) if len(lines) > 1 else ""
                
                slides.append({
                    "slide_number": page_num,
                    "title": title,
                    "content": content
                })
    except Exception as e:
        print(f"Error extracting PDF content: {e}")
    
    return slides

def parse_timestamp(timestamp_str):
    """Parse timestamp string into seconds."""
    # Extract just the first timestamp from formats like "206\n00:23:30.099 --> 00:23:36.720"
    match = re.search(r'(\d{2}):(\d{2}):(\d{2})\.(\d{3})', timestamp_str)
    if not match:
        return 0
    
    hours, minutes, seconds, milliseconds = map(int, match.groups())
    total_seconds = hours * 3600 + minutes * 60 + seconds + milliseconds / 1000
    return total_seconds

def extract_raw_transcript(transcript_path):
    """Extract transcript lines with timestamps from transcript file."""
    with open(transcript_path, 'r', encoding='utf-8') as f:
        lines = f.readlines()
    
    raw_lines = []
    current_timestamp = None
    current_text = ""
    
    for line in lines:
        line = line.strip()
        
        # Check for timestamp lines (e.g., "206\n00:23:30.099 --> 00:23:36.720")
        timestamp_match = re.search(r'\d{2}:\d{2}:\d{2}\.\d{3} -->', line)
        if timestamp_match:
            current_timestamp = line
            continue
            
        # Check for speaker line (e.g., "James Wagner: ...")
        speaker_match = re.match(r'^([^:]+):\s*(.+)$', line)
        if speaker_match and current_timestamp:
            speaker = speaker_match.group(1).strip()
            text = speaker_match.group(2).strip()
            
            if text:  # Only add non-empty lines
                raw_lines.append({
                    "timestamp": current_timestamp,
                    "speaker": speaker,
                    "text": text,
                    "seconds": parse_timestamp(current_timestamp)
                })
            
            current_timestamp = None
    
    # Sort by timestamp (seconds)
    raw_lines.sort(key=lambda x: x["seconds"])
    return raw_lines

def preprocess_text(text):
    """Clean text for better matching."""
    if not text:
        return ""
    
    # Convert to lowercase
    text = text.lower()
    
    # Remove special characters but keep spaces between words
    text = re.sub(r'[^\w\s]', ' ', text)
    
    # Remove extra whitespace
    text = re.sub(r'\s+', ' ', text).strip()
    
    return text

def match_transcript_to_slides(transcript_lines, slides, similarity_threshold=0.5):
    """Match transcript lines to slides using TF-IDF similarity with a threshold."""
    # Preprocess slide content for matching
    slide_texts = []
    for slide in slides:
        # Combine title and content with higher weight for title
        slide_text = f"{slide['title']} {slide['title']} {slide['content']}"
        slide_texts.append(preprocess_text(slide_text))
    
    # Preprocess transcript lines
    transcript_texts = [preprocess_text(line["text"]) for line in transcript_lines]
    
    # Create TF-IDF vectors
    vectorizer = TfidfVectorizer(stop_words='english', min_df=2, max_df=0.95)
    
    # Fit vectorizer on all texts to build vocabulary
    all_texts = slide_texts + transcript_texts
    vectorizer.fit(all_texts)
    
    # Transform slides and transcript lines to TF-IDF vectors
    slide_vectors = vectorizer.transform(slide_texts)
    transcript_vectors = vectorizer.transform(transcript_texts)
    
    # Calculate similarity between each transcript line and each slide
    similarities = cosine_similarity(transcript_vectors, slide_vectors)
    
    # Find the best matching slide for each transcript line
    matches = []
    unmatched_count = 0
    
    for i, line in enumerate(transcript_lines):
        # Get the slide with highest similarity
        best_slide_idx = np.argmax(similarities[i])
        similarity_score = similarities[i][best_slide_idx]
        
        # Only assign to a slide if similarity exceeds threshold
        if similarity_score >= similarity_threshold:
            matches.append({
                "timestamp": line["timestamp"],
                "seconds": line["seconds"],
                "speaker": line.get("speaker", ""),
                "transcript_line": line["text"],
                "text": line["text"],  # Added for compatibility with generate_transcript_by_slide
                "slide_number": slides[best_slide_idx]["slide_number"],
                "slide_title": slides[best_slide_idx]["title"],
                "confidence": float(similarity_score)
            })
        else:
            # Mark as unmatched with a special slide_number
            unmatched_count += 1
            matches.append({
                "timestamp": line["timestamp"],
                "seconds": line["seconds"],
                "speaker": line.get("speaker", ""),
                "transcript_line": line["text"],
                "text": line["text"],
                "slide_number": 0,  # Use 0 to indicate unmatched
                "slide_title": "Unmatched",
                "confidence": float(similarity_score)
            })
    
    if unmatched_count > 0:
        print(f"Note: {unmatched_count} transcript lines did not meet the similarity threshold and were marked as unmatched.")
    
    return matches

def apply_temporal_coherence(matches, slides, window_size=5):
    """
    Post-processing step to enforce temporal coherence.
    We expect slides to progress forward in time, with occasional backwards references.
    """
    if not matches:
        return []
    
    processed_matches = []
    current_slide = matches[0]["slide_number"]
    
    for i, match in enumerate(matches):
        # For low confidence matches, look at neighboring matches
        if match["confidence"] < 0.1:  # Threshold can be adjusted
            # Get a window of matches before and after this match
            start_idx = max(0, i - window_size)
            end_idx = min(len(matches), i + window_size + 1)
            window = matches[start_idx:end_idx]
            
            # Count slides in window to find most common
            slide_counts = {}
            for m in window:
                slide_num = m["slide_number"]
                if slide_num not in slide_counts:
                    slide_counts[slide_num] = 0
                slide_counts[slide_num] += 1
            
            # Find most common slide in window
            most_common_slide = max(slide_counts.items(), key=lambda x: x[1])[0]
            
            # Override low confidence match with most common slide
            match["slide_number"] = most_common_slide
            match["original_confidence"] = match["confidence"]
            match["confidence"] = -1  # Mark as corrected
            
            # Update slide title
            for slide in slides:
                if slide["slide_number"] == most_common_slide:
                    match["slide_title"] = slide["title"]
                    break
        
        processed_matches.append(match)
    
    return processed_matches

def generate_transcript_by_slide(mapping, slides, output_path):
    """
    Generate a formatted text file with transcript lines grouped by slide.
    
    Args:
        mapping: List of transcript lines mapped to slides
        slides: List of slide information
        output_path: Path to save the formatted output
    """
    # Group transcript lines by slide number
    slides_dict = {slide["slide_number"]: slide for slide in slides}
    transcript_by_slide = {}
    
    for match in mapping:
        slide_num = match["slide_number"]
        
        # Skip unmatched transcript lines (slide_number = 0)
        if slide_num == 0:
            continue
            
        if slide_num not in transcript_by_slide:
            transcript_by_slide[slide_num] = []
        
        # Format timestamp for readability
        timestamp = match["timestamp"]
        time_match = re.search(r'(\d{2}:\d{2}:\d{2})', timestamp)
        formatted_time = time_match.group(1) if time_match else "00:00:00"
        
        # Add formatted transcript line
        transcript_by_slide[slide_num].append({
            "time": formatted_time,
            "speaker": match.get("speaker", "Speaker"),
            "text": match["text"]
        })
    
    # Write formatted output
    with open(output_path, "w", encoding="utf-8") as f:
        f.write("TRANSCRIPT BY SLIDE\n")
        f.write("===================\n\n")
        
        # First create a section for unmatched lines if there are any
        unmatched = [m for m in mapping if m["slide_number"] == 0]
        if unmatched:
            f.write("UNMATCHED TRANSCRIPT LINES\n")
            f.write("-" * 80 + "\n")
            for line in unmatched:
                time_match = re.search(r'(\d{2}:\d{2}:\d{2})', line["timestamp"])
                formatted_time = time_match.group(1) if time_match else "00:00:00"
                f.write(f"[{formatted_time}] {line.get('speaker', 'Speaker')}: {line['text']}\n")
            f.write("\n\n")
        
        # Then write all slides with transcripts
        for slide_num in sorted(transcript_by_slide.keys()):
            # Get slide info
            slide_info = slides_dict.get(slide_num, {"title": f"Slide {slide_num}", "content": ""})
            slide_title = slide_info["title"] or f"Slide {slide_num}"
            
            # Write slide header
            f.write(f"SLIDE {slide_num}: {slide_title}\n")
            f.write("-" * 80 + "\n")
            
            # Write all transcript lines for this slide
            for line in transcript_by_slide[slide_num]:
                f.write(f"[{line['time']}] {line['speaker']}: {line['text']}\n")
            
            # Add empty line between slides
            f.write("\n\n")
    
    return len(transcript_by_slide), len(unmatched)

def main():
    # File paths
    pptx_path = "./Lecture01_Example.pptx"
    transcript_path = "./transcript_vtt.vtt"
    # json_output_path = "v3_mapping.json"
    txt_output_path = "output.txt"
    
    # Configurable parameters
    similarity_threshold = 0.05  # Minimum similarity score to assign a transcript line to a slide
    
    # 1. Extract slides (prioritize PPTX if available)
    print(f"Extracting slides from PowerPoint: {pptx_path}")
    slides = extract_slides_from_pptx(pptx_path)
    
    
    # 3. Extract transcript
    if not os.path.exists(transcript_path):
        raise FileNotFoundError(f"Transcript file not found: {transcript_path}")
    
    raw_transcript = extract_raw_transcript(transcript_path)
    print(f"Extracted {len(raw_transcript)} transcript lines")
    
    # 4. Match transcript to slides using TF-IDF similarity with threshold
    print(f"Matching transcript lines to slides (similarity threshold: {similarity_threshold})...")
    mapping = match_transcript_to_slides(raw_transcript, slides, similarity_threshold)
    
    # 5. Apply temporal coherence as post-processing
    print("Applying temporal coherence...")
    mapping = apply_temporal_coherence(mapping, slides)
    
    
    # 7. Generate formatted transcript by slide
    print("Generating transcript by slide...")
    generate_transcript_by_slide(mapping, slides, txt_output_path)
    print(f"Formatted transcript saved to {txt_output_path}")
    
    # 8. Print summary
    matched_lines = [m for m in mapping if m["slide_number"] > 0]
    unmatched_lines = [m for m in mapping if m["slide_number"] == 0]
    
    print("\nSummary:")
    print(f"Total transcript lines: {len(mapping)}")
    print(f"Matched transcript lines: {len(matched_lines)} ({len(matched_lines)/len(mapping)*100:.1f}%)")
    print(f"Unmatched transcript lines: {len(unmatched_lines)} ({len(unmatched_lines)/len(mapping)*100:.1f}%)")
    
    print("\nTranscript lines per slide:")
    slide_counts = {}
    for match in matched_lines:
        slide_num = match["slide_number"]
        if slide_num not in slide_counts:
            slide_counts[slide_num] = 0
        slide_counts[slide_num] += 1
    
    for slide_num in sorted(slide_counts.keys()):
        print(f"Slide {slide_num}: {slide_counts[slide_num]} lines")

if __name__ == "__main__":
    main()
