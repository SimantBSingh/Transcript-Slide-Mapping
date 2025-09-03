#!/usr/bin/env python3
import os
import re
import io
import pickle
from dataclasses import dataclass
from typing import List, Optional
from collections import defaultdict, deque
from datetime import datetime

import numpy as np
from rank_bm25 import BM25Okapi
import gensim
from gensim.models import Word2Vec
from sklearn.metrics.pairwise import cosine_similarity
from pptx import Presentation
from pptx.enum.shapes import PP_PLACEHOLDER
import nltk
from nltk.stem import PorterStemmer
from nltk.corpus import stopwords

# --------------------
# Data Models
# --------------------
@dataclass
class Utterance:
    speaker: str
    text: str
    start_sec: float
    end_sec: float
    mid_sec: float
    duration: float
    utterance_id: Optional[int] = None
    slide_number: int = 0
    bm25_score: float = 0.0
    w2v_score: float = 0.0
    combined_score: float = 0.0

@dataclass
class Slide:
    slide_number: int
    title: str
    raw_text: str
    content: str

# --------------------
# Setup
# --------------------
nltk.download('punkt', quiet=True)
nltk.download('stopwords', quiet=True)
stemmer = PorterStemmer()
STOP_WORDS = set(stopwords.words('english'))

# --------------------
# Utility Functions
# --------------------
def preprocess(text: str) -> List[str]:
    # tokenize, lowercase, remove stopwords, stem
    tokens = nltk.word_tokenize(text)
    processed = []
    for tok in tokens:
        tok = tok.lower()
        if tok.isalpha() and tok not in STOP_WORDS:
            processed.append(stemmer.stem(tok))
    return processed

# --------------------
# Extract Slides
# --------------------
def extract_slides(pptx_path: str) -> List[Slide]:
    prs = Presentation(pptx_path)
    slides = []
    for idx, sl in enumerate(prs.slides, start=1):
        title = ''
        placeholders = [s for s in sl.shapes if s.is_placeholder and s.placeholder_format.type == PP_PLACEHOLDER.TITLE]
        if placeholders:
            shp = placeholders[0]
            if shp.has_text_frame:
                title = shp.text.strip()
        if not title:
            for s in sl.shapes:
                if s.has_text_frame and s.text.strip():
                    title = s.text.strip()
                    break
        texts = []
        for s in sl.shapes:
            if s.has_text_frame:
                for p in s.text_frame.paragraphs:
                    t = p.text.strip()
                    if t:
                        texts.append(t)
        rawtxt = title + ' ' + ' '.join(texts)
        slides.append(Slide(idx, title, rawtxt, ' '.join(texts)))
    return slides

# --------------------
# Extract Transcript
# --------------------
def parse_timestamp(ts: str) -> dict:
    m = re.findall(r"(\d+):(\d+):(\d+)\.(\d+)", ts)
    if len(m) != 2:
        return {'start_sec':0,'end_sec':0,'mid_sec':0,'duration':0}
    def to_sec(h,m,s,ms): return int(h)*3600+int(m)*60+int(s)+int(ms)/1000
    start = to_sec(*m[0]); end = to_sec(*m[1])
    return {'start_sec':start,'end_sec':end,'mid_sec':(start+end)/2,'duration':end-start}

from typing import List

def extract_raw_transcript(vtt_path: str) -> List[Utterance]:
    lines = [l.strip() for l in open(vtt_path, encoding='utf-8')]
    raw = []
    curr_ts = None
    curr_spk = 'Unknown'
    buf: List[str] = []
    for ln in lines:
        if '-->' in ln:
            curr_ts = parse_timestamp(ln)
            continue
        m = re.match(r'^([A-Za-z ]+):\s(.+)$', ln)
        if m:
            if buf and curr_ts:
                text = ' '.join(buf)
                # Strip trailing numeric codes
                text = re.sub(r'\s*\d+\s*$', '', text)
                raw.append(Utterance(curr_spk, text,
                                     curr_ts['start_sec'], curr_ts['end_sec'],
                                     curr_ts['mid_sec'], curr_ts['duration']))
                buf = []
            curr_spk = m.group(1).strip()
            buf.append(m.group(2).strip())
        elif curr_ts:
            buf.append(ln)
    if buf and curr_ts:
        text = ' '.join(buf)
        text = re.sub(r'\s*\d+\s*$', '', text)
        raw.append(Utterance(curr_spk, text,
                             curr_ts['start_sec'], curr_ts['end_sec'],
                             curr_ts['mid_sec'], curr_ts['duration']))
    raw.sort(key=lambda u: u.start_sec)
    for i, u in enumerate(raw, start=1): u.utterance_id = i
    return raw

# --------------------
# Build BM25 and Word2Vec Models
# --------------------
def build_models(slides: List[Slide], transcript: List[Utterance]):
    # Prepare corpus
    slide_docs = [preprocess(s.raw_text) for s in slides]
    utterance_docs = [preprocess(u.text) for u in transcript]
    # BM25
    bm25 = BM25Okapi(slide_docs)
    # Word2Vec
    combined_docs = slide_docs + utterance_docs
    w2v_model = Word2Vec(sentences=combined_docs, vector_size=100, window=5, min_count=1, workers=4)
    return bm25, w2v_model

# --------------------
# Matching Function
# --------------------
def match_with_bm25_w2v(slides: List[Slide], utts: List[Utterance],
                        bm25: BM25Okapi, w2v_model: Word2Vec,
                        alpha: float = 0.5) -> List[Utterance]:
    slide_vecs = []
    # Precompute slide average embeddings
    for s in slides:
        toks = preprocess(s.raw_text)
        vecs = [w2v_model.wv[t] for t in toks if t in w2v_model.wv]
        slide_vecs.append(np.mean(vecs, axis=0) if vecs else np.zeros(w2v_model.vector_size))
    # Match each utterance
    for u in utts:
        toks = preprocess(u.text)
        # BM25 score on raw tokens
        bm25_scores = bm25.get_scores(toks)
        best_idx = int(np.argmax(bm25_scores))
        bm25_best = float(bm25_scores[best_idx])
        # W2V embedding similarity
        utt_vecs = [w2v_model.wv[t] for t in toks if t in w2v_model.wv]
        utt_vec = np.mean(utt_vecs, axis=0) if utt_vecs else np.zeros(w2v_model.vector_size)
        sims = cosine_similarity([utt_vec], slide_vecs)[0]
        w2v_best_idx = int(np.argmax(sims))
        w2v_best = float(sims[w2v_best_idx])
        # Combine
        combined = alpha * (bm25_best / (bm25_best + 1e-6)) + (1-alpha) * w2v_best
        # Choose slide: if indices match, pick that, else max combined among both
        if best_idx == w2v_best_idx:
            chosen = best_idx
        else:
            # Try both candidates
            cand_scores = {best_idx: combined, w2v_best_idx: combined}
            chosen = max(cand_scores, key=cand_scores.get)
        u.slide_number = slides[chosen].slide_number
        u.bm25_score = bm25_best
        u.w2v_score = w2v_best
        u.combined_score = combined
    return utts

# --------------------
# Output Generation
# --------------------
def format_time(sec: float) -> str:
    return datetime.utcfromtimestamp(sec).strftime('%H:%M:%S')


def generate_output_by_slide(matches: List[Utterance], slides: List[Slide], out_path: str):
    groups = defaultdict(list)
    for m in matches:
        groups[m.slide_number].append(m)
    meta = {s.slide_number: s for s in slides}
    with open(out_path, 'w', encoding='utf-8') as f:
        f.write("TRANSCRIPT BY SLIDE\n==================\n\n")
        for sl_num in sorted(groups):
            if sl_num == 0: continue
            title = meta[sl_num].title or f"Slide {sl_num}"
            f.write(f"SLIDE {sl_num}: {title}\n")
            f.write('-'*80 + '\n')
            for u in groups[sl_num]:
                f.write(f"[{format_time(u.start_sec)}] {u.speaker}: {u.text}\n")
            f.write('\n')

# --------------------
# Main Entry
# --------------------
def main(transcript_path: str, pptx_path: str, output_path: str):
    print("Extracting Slides...\n")
    slides = extract_slides(pptx_path)
    print("Extracting Transcript...\n")
    utts = extract_raw_transcript(transcript_path)
    print("Building bm25 and w2v models...\n")
    bm25, w2v_model = build_models(slides, utts)
    print("Matching Transcript Line with slides...\n")
    matches = match_with_bm25_w2v(slides, utts, bm25, w2v_model)
    print("Generating Outout textfile...\n")
    generate_output_by_slide(matches, slides, output_path)
    print(f"Output written to {output_path}")

if __name__ == '__main__':
    main('transcript.txt', 'Lecture01_Example.pptx', 'v3_output.txt')
